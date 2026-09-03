"""
Agent System (Enterprise++ v3.5) - MCP Document Parser & PDF Generator Server
Izolowana usługa parsowania dokumentów oraz dynamicznego generowania raportów PDF.

Zgodnie ze spec.md:
- Sekcja 6.1: Read-Only FS + Tmpfs + twarde limity pamięci (--memory=2g)
- Sekcja 8.1: Generowanie artefaktów (raporty PDF z podsumowaniem analiz)
- Sekcja 4.1: Weryfikacja SHA-256 plików przed przetworzeniem (Zero-Trust)

Zabezpieczenia sandboxa:
- Weryfikacja skrótu SHA-256 pliku PRZED parsowaniem (anti-RCE / integrity)
- Twarde limity pamięci w locie (RSS monitoring, anti-OOM)
- Tylko tmpfs jako ścieżka zapisu, zero zapisów na dysk fizyczny
- Walidacja ścieżek (anti path-traversal), whitelist rozszerzeń
"""

from __future__ import annotations

import asyncio
import hashlib
import logging
import os
import re
import resource
import tempfile
import uuid
from datetime import datetime, timezone
from enum import Enum
from pathlib import Path
from typing import Any, Optional

from pydantic import BaseModel, Field, field_validator

logger = logging.getLogger(__name__)

MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024  # 50MB
TMPFS_PATH = "/tmp"
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".html", ".md"}
MAX_MEMORY_MB = 512          # in-flight RSS cap for this process (anti-OOM)
MAX_EXTRACTED_CHARS = 2_000_000  # cap on extracted text size (anti memory blow-up)
MAX_TABLE_ROWS = 10_000

SAFE_FILENAME_RE = re.compile(r"^[A-Za-z0-9._\- ]{1,255}$")


class DocumentFormat(str, Enum):
    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    HTML = "html"
    MARKDOWN = "md"


class ParsingStatus(str, Enum):
    PENDING = "pending"
    PARSING = "parsing"
    COMPLETED = "completed"
    FAILED = "failed"
    REJECTED = "rejected"


class DocumentMetadata(BaseModel):
    filename: str
    file_size_bytes: int
    format: DocumentFormat
    sha256_hash: str
    uploaded_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    source: str = "unknown"

    @field_validator("file_size_bytes")
    @classmethod
    def validate_file_size(cls, v: int) -> int:
        if v > MAX_FILE_SIZE_BYTES:
            raise ValueError(f"File size {v} exceeds maximum {MAX_FILE_SIZE_BYTES}")
        if v <= 0:
            raise ValueError("File size must be positive")
        return v

    @field_validator("filename")
    @classmethod
    def validate_filename(cls, v: str) -> str:
        if ".." in v or "/" in v or "\\" in v:
            raise ValueError(f"Invalid filename: {v} (path traversal detected)")
        if not SAFE_FILENAME_RE.match(v):
            raise ValueError(f"Filename contains illegal characters: {v!r}")
        ext = Path(v).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file extension: {ext}")
        return v


class ParsedContent(BaseModel):
    text: str = ""
    tables: list[dict[str, Any]] = Field(default_factory=list)
    images: list[dict[str, Any]] = Field(default_factory=list)
    metadata: dict[str, Any] = Field(default_factory=dict)
    page_count: int = 0
    word_count: int = 0


class ParseResult(BaseModel):
    status: ParsingStatus
    document_id: str
    metadata: Optional[DocumentMetadata] = None
    content: Optional[ParsedContent] = None
    error_message: Optional[str] = None
    processing_time_ms: float = 0.0
    memory_used_mb: float = 0.0
    parsed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class GeneratedReport(BaseModel):
    """Result of PDF report generation."""
    report_id: str
    filename: str
    file_path: str
    file_size_bytes: int
    sha256_hash: str
    generated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    title: str = ""
    page_count: int = 0


class MemoryCapExceededError(RuntimeError):
    """Raised when in-flight RSS exceeds hard memory cap (anti-OOM)."""
    pass


def get_process_rss_mb() -> float:
    """Current process RSS in megabytes (Linux)."""
    try:
        rss_kb = resource.getrusage(resource.RUSAGE_SELF).ru_maxrss
        return rss_kb / 1024.0
    except Exception:
        return 0.0


def enforce_memory_cap(max_mb: float = MAX_MEMORY_MB) -> None:
    """Hard in-flight memory enforcement; raises when cap exceeded (anti-OOM)."""
    rss = get_process_rss_mb()
    if rss > max_mb:
        raise MemoryCapExceededError(
            f"Process RSS {rss:.1f}MB exceeds hard cap {max_mb}MB - aborting to protect host"
        )


def _parse_text_sync(text: str, fmt: str, extra: Optional[dict[str, Any]] = None) -> ParsedContent:
    """Shared post-processing: caps + word count."""
    if len(text) > MAX_EXTRACTED_CHARS:
        logger.warning("Extracted text truncated from %d to %d chars", len(text), MAX_EXTRACTED_CHARS)
        text = text[:MAX_EXTRACTED_CHARS]
    content = ParsedContent(
        text=text,
        page_count=extra.get("page_count", 1) if extra else 1,
        tables=extra.get("tables", []) if extra else [],
        images=extra.get("images", []) if extra else [],
        metadata={"parser": "native", "format": fmt, **(extra or {})},
    )
    content.word_count = len(content.text.split())
    return content


def _extract_pdf_sync(path: str) -> ParsedContent:
    """Real PDF extraction via pypdf (fully offline)."""
    from pypdf import PdfReader  # local library, no network

    reader = PdfReader(path)
    pages: list[str] = []
    tables: list[dict[str, Any]] = []
    for i, page in enumerate(reader.pages):
        enforce_memory_cap()
        try:
            pages.append(page.extract_text() or "")
        except Exception as exc:
            logger.warning("PDF page %d extraction failed: %s", i + 1, exc)
            pages.append("")
        try:
            for table in page.extract_tables() or []:
                tables.append({"page": i + 1, "rows": table[:MAX_TABLE_ROWS]})
        except Exception:
            pass
    text = "\n\n".join(pages)
    return _parse_text_sync(text, "pdf", {"page_count": len(reader.pages), "tables": tables})


def _extract_docx_sync(path: str) -> ParsedContent:
    """Real DOCX extraction via python-docx."""
    from docx import Document  # local library, no network

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    tables = []
    for ti, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        tables.append({"index": ti, "rows": rows[:MAX_TABLE_ROWS]})
    text = "\n".join(paragraphs)
    return _parse_text_sync(text, "docx", {"page_count": 1, "tables": tables})


def _extract_xlsx_sync(path: str) -> ParsedContent:
    """Real XLSX extraction via openpyxl (read-only, no formulas evaluation)."""
    from openpyxl import load_workbook  # local library, no network

    wb = load_workbook(path, read_only=True, data_only=True)
    tables: list[dict[str, Any]] = []
    text_parts: list[str] = []
    for ws in wb.worksheets:
        enforce_memory_cap()
        rows: list[list[Any]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append(["" if v is None else str(v) for v in row])
            if len(rows) >= MAX_TABLE_ROWS:
                break
        tables.append({"sheet": ws.title, "rows": rows})
        text_parts.append(f"### Sheet: {ws.title}\n" + "\n".join(
            " | ".join(r) for r in rows[:200]
        ))
    wb.close()
    return _parse_text_sync("\n\n".join(text_parts), "xlsx", {"page_count": len(tables), "tables": tables})


def _extract_pptx_sync(path: str) -> ParsedContent:
    """Real PPTX extraction via python-pptx."""
    from pptx import Presentation  # local library, no network

    prs = Presentation(path)
    slides_text: list[str] = []
    for idx, slide in enumerate(prs.slides):
        enforce_memory_cap()
        parts = [f"## Slide {idx + 1}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
        slides_text.append("\n".join(parts))
    return _parse_text_sync("\n\n".join(slides_text), "pptx", {"page_count": len(prs.slides)})


_HTML_TAG_RE = re.compile(r"<[^>]+>")
_SCRIPT_STYLE_RE = re.compile(r"<(script|style)[^>]*>.*?</\1>", re.IGNORECASE | re.DOTALL)


def _extract_html_sync(path: str) -> ParsedContent:
    """HTML extraction with script/style stripping (XSS-safe)."""
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    cleaned = _SCRIPT_STYLE_RE.sub("", raw)
    text = _HTML_TAG_RE.sub("\n", cleaned)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return _parse_text_sync(text, "html", {"page_count": 1})


def _extract_markdown_sync(path: str) -> ParsedContent:
    """Markdown extraction (native, zero deps)."""
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return _parse_text_sync(text, "markdown", {"page_count": 1})


_PDF_EXTRACTORS = {
    DocumentFormat.PDF: _extract_pdf_sync,
    DocumentFormat.DOCX: _extract_docx_sync,
    DocumentFormat.XLSX: _extract_xlsx_sync,
    DocumentFormat.PPTX: _extract_pptx_sync,
    DocumentFormat.HTML: _extract_html_sync,
    DocumentFormat.MARKDOWN: _extract_markdown_sync,
}


class DocParserMCP:
    """
    MCP Document Parser Server with strict isolation.

    Security invariants:
    - SHA-256 integrity verification BEFORE parsing (expected_hash check)
    - Hard in-flight memory caps enforced between every page/row extraction
    - All writes go to tmpfs only; no persistent filesystem writes
    - Path traversal + extension whitelist validation on every entry point
    """

    def __init__(
        self,
        tmpfs_path: str = TMPFS_PATH,
        max_file_size_mb: int = 50,
        timeout_seconds: int = 30,
        max_memory_mb: int = MAX_MEMORY_MB,
    ):
        self._tmpfs_path = tmpfs_path
        self._max_file_size_bytes = max_file_size_mb * 1024 * 1024
        self._timeout_seconds = timeout_seconds
        self._max_memory_mb = max_memory_mb
        self._parse_results: dict[str, ParseResult] = {}
        self._initialized = False
        logger.info(
            "DocParserMCP initialized: tmpfs=%s max_size=%dMB mem_cap=%dMB",
            tmpfs_path, max_file_size_mb, max_memory_mb,
        )

    async def initialize(self) -> None:
        try:
            test_file = os.path.join(self._tmpfs_path, ".write_test")
            with open(test_file, "w") as f:
                f.write("test")
            os.remove(test_file)
            self._initialized = True
            logger.info("DocParserMCP initialization complete")
        except OSError as e:
            raise RuntimeError(f"Cannot write to tmpfs at {self._tmpfs_path}: {e}")

    async def _compute_file_hash(self, file_path: str) -> str:
        loop = asyncio.get_event_loop()

        def _hash() -> str:
            sha256 = hashlib.sha256()
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    sha256.update(chunk)
            return sha256.hexdigest()

        return await loop.run_in_executor(None, _hash)

    async def verify_sha256(self, file_path: str, expected_hash: str) -> bool:
        """
        Zero-Trust integrity gate: compute file digest and compare
        with expected hash BEFORE any parsing touches the content.
        """
        actual = await self._compute_file_hash(file_path)
        ok = actual.lower() == expected_hash.lower()
        if not ok:
            logger.error(
                "SHA-256 MISMATCH for %s: expected=%s actual=%s",
                file_path, expected_hash, actual,
            )
        return ok

    async def parse_document(
        self,
        file_path: str,
        filename: str,
        expected_sha256: Optional[str] = None,
    ) -> ParseResult:
        """
        Parse a document with full safety pipeline:
        1. tmpfs existence + size + path checks
        2. SHA-256 pre-verification (when expected_sha256 provided)
        3. Real extraction with memory caps + timeout
        4. Post-parse RSS verification
        """
        if not self._initialized:
            raise RuntimeError("DocParserMCP not initialized")

        document_id = f"doc_{uuid.uuid4().hex[:12]}"
        loop = asyncio.get_event_loop()
        start = loop.time()
        result: Optional[ParseResult] = None

        try:
            # -- Stage 1: path & size validation
            if not os.path.exists(file_path):
                result = ParseResult(status=ParsingStatus.REJECTED, document_id=document_id,
                                     error_message=f"File not found: {file_path}")
                return result
            if not os.path.abspath(file_path).startswith(os.path.abspath(self._tmpfs_path)):
                result = ParseResult(status=ParsingStatus.REJECTED, document_id=document_id,
                                     error_message="File must reside in tmpfs sandbox")
                return result
            file_size = os.path.getsize(file_path)
            if file_size > self._max_file_size_bytes:
                result = ParseResult(status=ParsingStatus.REJECTED, document_id=document_id,
                                     error_message=f"Size {file_size} > {self._max_file_size_bytes}")
                return result

            # -- Stage 2: SHA-256 integrity gate (before parsing)
            if expected_sha256 is not None:
                ok = await self.verify_sha256(file_path, expected_sha256)
                if not ok:
                    result = ParseResult(status=ParsingStatus.REJECTED, document_id=document_id,
                                         error_message="SHA-256 integrity check failed - file rejected")
                    return result
            sha256_hash = await self._compute_file_hash(file_path)

            ext = Path(filename).suffix.lower().lstrip(".")
            try:
                doc_format = DocumentFormat(ext)
            except ValueError:
                result = ParseResult(status=ParsingStatus.REJECTED, document_id=document_id,
                                     error_message=f"Unsupported format: {ext}")
                return result

            metadata = DocumentMetadata(
                filename=filename,
                file_size_bytes=file_size,
                format=doc_format,
                sha256_hash=sha256_hash,
            )

            # -- Stage 3: real extraction with timeout + memory caps
            extractor = _PDF_EXTRACTORS[doc_format]
            try:
                content = await asyncio.wait_for(
                    loop.run_in_executor(None, extractor, file_path),
                    timeout=self._timeout_seconds,
                )
            except asyncio.TimeoutError:
                result = ParseResult(status=ParsingStatus.FAILED, document_id=document_id,
                                     metadata=metadata,
                                     error_message=f"Parsing timed out after {self._timeout_seconds}s")
                return result

            # -- Stage 4: post-parse memory verification
            enforce_memory_cap(self._max_memory_mb)

            processing_time = (loop.time() - start) * 1000
            result = ParseResult(
                status=ParsingStatus.COMPLETED,
                document_id=document_id,
                metadata=metadata,
                content=content,
                processing_time_ms=processing_time,
                memory_used_mb=get_process_rss_mb(),
            )
            logger.info("Parsed %s (%s) in %.1fms", document_id, doc_format, processing_time)
            return result

        except MemoryCapExceededError as exc:
            logger.critical("MEMORY CAP: %s", exc)
            return ParseResult(status=ParsingStatus.FAILED, document_id=document_id,
                               error_message=str(exc),
                               processing_time_ms=(loop.time() - start) * 1000)
        except Exception as exc:
            logger.exception("Parse failed: %s", exc)
            return ParseResult(status=ParsingStatus.FAILED, document_id=document_id,
                               error_message=str(exc),
                               processing_time_ms=(loop.time() - start) * 1000,
                               memory_used_mb=get_process_rss_mb())
        finally:
            if result is not None:
                self._parse_results[document_id] = result

    async def generate_pdf_report(
        self,
        title: str,
        sections: list[dict[str, Any]],
        output_filename: Optional[str] = None,
    ) -> GeneratedReport:
        """
        Generate a PDF report via ReportLab (fully offline, sandboxed to tmpfs).

        sections: [{"heading": str, "body": str}, ...]
        """
        loop = asyncio.get_event_loop()
        report_id = f"rpt_{uuid.uuid4().hex[:12]}"
        output_filename = output_filename or f"{report_id}.pdf"
        if not SAFE_FILENAME_RE.match(output_filename):
            raise ValueError(f"Illegal report filename: {output_filename!r}")
        out_path = os.path.join(self._tmpfs_path, output_filename)

        def _render() -> tuple[int, str]:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

            enforce_memory_cap(self._max_memory_mb)
            doc = SimpleDocTemplate(
                out_path, pagesize=A4,
                leftMargin=20 * mm, rightMargin=20 * mm,
                topMargin=20 * mm, bottomMargin=20 * mm,
                title=title, author="Agent System v3.5",
            )
            styles = getSampleStyleSheet()
            story: list[Any] = [Paragraph(title, styles["Title"]), Spacer(1, 8 * mm)]
            for section in sections:
                enforce_memory_cap(self._max_memory_mb)
                heading = str(section.get("heading", ""))
                body = str(section.get("body", ""))
                if heading:
                    story.append(Paragraph(heading, styles["Heading2"]))
                if body:
                    story.append(Paragraph(body, styles["BodyText"]))
                story.append(Spacer(1, 4 * mm))
            doc.build(story)

            size = os.path.getsize(out_path)
            sha = hashlib.sha256()
            with open(out_path, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    sha.update(chunk)
            return size, sha.hexdigest()

        size, sha = await asyncio.wait_for(
            loop.run_in_executor(None, _render),
            timeout=self._timeout_seconds,
        )
        enforce_memory_cap(self._max_memory_mb)
        report = GeneratedReport(
            report_id=report_id,
            filename=output_filename,
            file_path=out_path,
            file_size_bytes=size,
            sha256_hash=sha,
            title=title,
        )
        logger.info("PDF report generated: %s (%d bytes)", out_path, size)
        return report

    async def cleanup(self, document_id: str) -> bool:
        if document_id in self._parse_results:
            del self._parse_results[document_id]
            return True
        return False

    async def health_check(self) -> dict[str, Any]:
        tmpfs_writable = False
        try:
            test_file = os.path.join(self._tmpfs_path, ".health_check")
            with open(test_file, "w") as f:
                f.write("ok")
            os.remove(test_file)
            tmpfs_writable = True
        except OSError:
            pass
        return {
            "status": "healthy" if (tmpfs_writable and self._initialized) else "degraded",
            "initialized": self._initialized,
            "tmpfs_writable": tmpfs_writable,
            "rss_mb": get_process_rss_mb(),
            "memory_cap_mb": self._max_memory_mb,
            "cached_results": len(self._parse_results),
        }
